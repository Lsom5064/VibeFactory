#!/usr/bin/env python3
"""Refresh the interaction-analysis workbook and build a Slides-compatible deck."""

from __future__ import annotations

import shutil
import unicodedata
from collections import Counter, OrderedDict
from copy import copy
from pathlib import Path
from typing import Iterable

from openpyxl import load_workbook
from openpyxl.styles import Alignment
from openpyxl.utils import get_column_letter
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
EXPORTS = ROOT / "exports"
WORKBOOK_PATH = EXPORTS / "participant_usage_history_google_sheets_numbered_images_and_analysis.xlsx"
BACKUP_PATH = EXPORTS / "participant_usage_history_google_sheets_numbered_images_and_analysis_before_20260810.xlsx"
IMAGES_DIR = EXPORTS / "participant_usage_history_numbered_images"
OUTPUT_PPTX = EXPORTS / "image_vs_text_interaction_analysis_20260810.pptx"


LABELS = ("성공", "일부 실패", "실패")
IMAGE_MODE = "이미지 첨부"
TEXT_MODE = "텍스트"
APP_STAGE = "앱 생성·수정 요청"


def find_reference_pptx() -> Path:
    for candidate in (Path.home() / "Downloads").glob("*.pptx"):
        normalized = unicodedata.normalize("NFC", candidate.name)
        if "제목 없는 프레젠테이션" in normalized:
            return candidate
    raise FileNotFoundError("다운로드 폴더에서 기본 참고 PPTX를 찾지 못했습니다.")


def clean(value: object) -> str:
    return str(value or "").strip()


def short_text(value: object, limit: int = 240) -> str:
    text = " ".join(clean(value).replace("\v", " ").split())
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip() + "…"


def label_slug(label: str) -> str:
    return {"성공": "success", "일부 실패": "partial", "실패": "failure"}[label]


def copy_row_style(ws, source_row: int, target_row: int, columns: int) -> None:
    for column in range(1, columns + 1):
        source = ws.cell(source_row, column)
        target = ws.cell(target_row, column)
        if source.has_style:
            target._style = copy(source._style)
        if source.number_format:
            target.number_format = source.number_format


def reset_sheet_rows(ws, rows: list[list[str]]) -> None:
    old_max = ws.max_row
    template_row = 2 if old_max >= 2 else 1
    required_max = len(rows) + 1
    for row_index, row_values in enumerate(rows, start=2):
        copy_row_style(ws, template_row, row_index, ws.max_column)
        for column, value in enumerate(row_values, start=1):
            ws.cell(row_index, column).value = value
            ws.cell(row_index, column).alignment = Alignment(
                horizontal="left",
                vertical="top",
                wrap_text=True,
            )
        ws.row_dimensions[row_index].height = 96 if ws.title == "Task별 방식 모음" else 74
    if old_max > required_max:
        ws.delete_rows(required_max + 1, old_max - required_max)


def interaction_rows(workbook) -> list[dict[str, object]]:
    source = workbook["참가자별 시간순 기록"]
    classified = workbook["상호작용 분류"]

    # The latest manual labels were entered in the chronological source sheet.
    # Mirror them into the analysis sheet using its source-row pointer.
    for row in range(2, classified.max_row + 1):
        try:
            source_row = int(classified.cell(row, 13).value)
        except (TypeError, ValueError):
            continue
        if source_row < 2 or source_row > source.max_row:
            continue
        label = clean(source.cell(source_row, 9).value)
        reason = clean(source.cell(source_row, 10).value)
        if label in LABELS:
            classified.cell(row, 9).value = label
            classified.cell(row, 11).value = reason

    headers = [clean(classified.cell(1, column).value) for column in range(1, classified.max_column + 1)]
    rows: list[dict[str, object]] = []
    for row in range(2, classified.max_row + 1):
        item = {headers[column - 1]: classified.cell(row, column).value for column in range(1, classified.max_column + 1)}
        if clean(item.get("성공 판정")) in LABELS:
            item["_sheet_row"] = row
            rows.append(item)
    return rows


def format_case(item: dict[str, object]) -> str:
    source_row = clean(item.get("원본 시트 행"))
    label = clean(item.get("성공 판정"))
    request = clean(item.get("참가자 요청")) or "(이미지만 첨부)"
    files = clean(item.get("첨부 파일명"))
    reason = clean(item.get("성공/실패 근거"))
    file_suffix = f" [첨부: {files}]" if files else ""
    return f"[원본 {source_row}행 · {label}]{file_suffix} {request}\n근거: {reason or '미기재'}"


def summarize_cases(items: Iterable[dict[str, object]], mode: str, label: str) -> str:
    values = [format_case(item) for item in items if clean(item.get("분석 방식")) == mode and clean(item.get("성공 판정")) == label]
    return "\n\n".join(values) if values else "해당 사례 없음"


def count_mode(items: Iterable[dict[str, object]], mode: str) -> str:
    counts = Counter(clean(item.get("성공 판정")) for item in items if clean(item.get("분석 방식")) == mode)
    return f"{counts['성공']} / {counts['일부 실패']} / {counts['실패']}"


def update_task_collection(workbook, rows: list[dict[str, object]]) -> None:
    ws = workbook["Task별 방식 모음"]
    old_transition = {clean(ws.cell(row, 3).value): clean(ws.cell(row, 11).value) for row in range(2, ws.max_row + 1)}
    grouped: OrderedDict[tuple[str, str, str], list[dict[str, object]]] = OrderedDict()
    for item in rows:
        key = (clean(item.get("참가자")), clean(item.get("Task 이름")), clean(item.get("Task ID")))
        grouped.setdefault(key, []).append(item)

    recent_reasons = {
        "c0adb5e9ffe749498263bc78dfee3223": (
            "달력 UI는 색상 표시·이동 위치·완성 배치를 구체적으로 적어 성공했다. 반면 알림 권한과 위젯은 화면 이미지만으로 원인을 확정하기 어려워 반복 실패했다. "
            "알림 문제는 상태 갱신·수집·복귀 조건을 번호로 정의한 요청에서 성공했고, 위젯은 표시 범위를 '오늘 일정'으로 축소한 뒤 성공했다."
        ),
        "a8cdbac1e46e4a1980dce71b62a59c88": (
            "이미지 없이도 년월·화살표·요일의 최종 배치를 구체적으로 적은 UI 요청은 성공했다. 다만 알림 권한 요청은 구체적이었어도 빌드 오류로 실패했고, '위젯을 다시 만들어줘'는 현재 실패 상태·기대 결과·보존 조건이 없어 실패했다. "
            "프롬프트 품질과 빌드·런타임 기술 실패를 분리해 평가해야 한다."
        ),
    }

    output: list[list[str]] = []
    for (participant, task_name, task_id), items in grouped.items():
        partial = [format_case(item) for item in items if clean(item.get("성공 판정")) == "일부 실패"]
        output.append([
            participant,
            task_name,
            task_id,
            count_mode(items, IMAGE_MODE),
            count_mode(items, TEXT_MODE),
            summarize_cases(items, IMAGE_MODE, "성공"),
            summarize_cases(items, TEXT_MODE, "성공"),
            summarize_cases(items, IMAGE_MODE, "실패"),
            summarize_cases(items, TEXT_MODE, "실패"),
            "\n\n".join(partial) if partial else "해당 사례 없음",
            recent_reasons.get(task_id) or old_transition.get(task_id) or "직접적인 실패→성공 연속 사례가 없어 추가 판단이 필요함.",
        ])
    reset_sheet_rows(ws, output)
    widths = [12, 24, 34, 18, 18, 52, 52, 52, 52, 52, 58]
    for index, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(index)].width = width
    ws.auto_filter.ref = f"A1:K{len(output) + 1}"
    ws.freeze_panes = "A2"


def style_analysis_sheet(ws, rows: list[list[str]], widths: list[int]) -> None:
    reset_sheet_rows(ws, rows)
    for index, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(index)].width = width
    ws.auto_filter.ref = f"A1:{get_column_letter(ws.max_column)}{len(rows) + 1}"
    ws.freeze_panes = "A2"


def update_characteristic_sheets(workbook, app_counts: dict[str, Counter]) -> None:
    image = app_counts[IMAGE_MODE]
    text = app_counts[TEXT_MODE]
    rate_text = (
        f"앱 생성·수정 요청 기준: 이미지 엄격 성공 {image['성공']}/{sum(image.values())}"
        f"({image['성공']/sum(image.values()):.1%}), 텍스트 {text['성공']}/{sum(text.values())}"
        f"({text['성공']/sum(text.values()):.1%}). 요청 난이도가 달라 인과 비교로 보지 않음."
    )

    success_rows = [
        ["이미지 첨부", "표시된 대상과 수정 동작을 함께 명시", "색·동그라미·번호로 화면 영역을 가리킨 뒤 삭제, 이동, 색상 변경처럼 한정된 동사를 사용했다.", "지도 6~10행, 감정진단 28~32행, 대화일정 144행", "이미지는 '어디'를, 텍스트는 '무엇을'을 확정할 때 가장 유용했다."],
        ["이미지 첨부", "최종 배치의 축·순서·위치를 구체화", "'1열 좌우', '가로 한 줄', '년월 왼쪽/오른쪽'처럼 완성 상태를 적었다.", "인바디코치 76행, 맞춤레시피 59행, 대화일정 144행", "상대적 표현보다 레이아웃 축과 요소 순서를 같이 제시할 때 오해가 줄었다."],
        ["이미지 첨부", "이미지·표시를 번호와 요구사항으로 1:1 대응", "각 이미지나 영역을 #1~#3으로 부르고 각각 별도 변경을 배정했다.", "인바디코치 76행, 출석체크90 136행", "참조 범위가 고정되어 여러 변경 중 누락이 줄었다."],
        ["이미지 첨부", "오류 화면과 재현 동작을 함께 제공", "오류가 보이는 화면, 누른 버튼, 기대 결과를 함께 적은 사례에서 성공했다.", "지도 4~6행", "단, 권한·위젯처럼 화면 밖 시스템 상태가 원인인 문제에서는 이미지만으로 부족했다."],
        ["이미지 첨부", "삭제 대상과 보존 대상의 경계를 명시", "부모 박스, 내부 텍스트, 버튼 중 무엇을 남길지 적었다.", "맞춤레시피 59행, 출석체크90 136행", "삭제와 보존을 함께 정의하면 인접 UI까지 지우는 오류를 줄일 수 있다."],
        ["텍스트", "트리거→상태 갱신→출력을 순서로 정의", "누구가 무엇을 했을 때 어떤 상태가 바뀌고 어디에 보여야 하는지 적었다.", "피아노진도장 86~90행, 대화일정 146행", "기능 요청은 화면 모양보다 상태 전이를 명시할 때 성공했다."],
        ["텍스트", "레이블·위치·순서를 정확하게 지정", "년월, 왼쪽/오른쪽 화살표, 요일 행의 최종 배치를 이미지 없이도 완전하게 적었다.", "대화일정 149행", "동일한 달력 배치가 이미지 144행과 텍스트 149행 모두에서 성공해 최종 상태 명시의 효과를 보여준다."],
        ["텍스트", "번호로 수용 기준을 나열", "권한 허용 후 즉시 갱신, 재실행 시 상태 복구, 실제 수집 등을 별도 항목으로 정의했다.", "대화일정 146행", "구현 단계와 검증 조건이 명확해졌다. 다만 150행은 동일한 요청이 빌드 오류로 실패해 표현과 기술 실패를 분리해야 한다."],
        ["텍스트", "외부 연동의 데이터 흐름과 경계를 정의", "외부 앱을 여는지, 데이터만 가져오는지, 어디에 저장·표시하는지를 분리했다.", "감정진단 33~36행", "'연동'이라는 단어보다 데이터 주체·이동·저장을 정의할 때 성공했다."],
        ["공통", "실패 후 범위를 줄이거나 폴백을 허용", "위젯 정보를 오늘·내일에서 오늘만으로 축소하고 규격을 줄여 성공했다.", "대화일정 145·147·151·154행", "위젯 연속 시도 4회 중 범위 축소·폴백을 적은 마지막 1회에서 성공했다."],
        ["공통", "요청 표현의 효과는 기술 성공과 분리해 해석", rate_text, "전체 앱 생성·수정 라벨 138건", "이미지 요청이 더 복잡한 문제에 사용됐을 수 있고, 좋은 요청도 빌드 오류로 실패할 수 있어 단순 비율을 인과로 보면 안 된다."],
    ]

    failure_rows = [
        ["이미지 첨부", "'이렇게', '다시'처럼 이미지 해석에 전적으로 의존", "수정 대상과 동작을 말하지 않고 이미지가 의도를 대신하게 했다.", "영양분석 108·110행", "보이는 요소가 여러 개면 어느 부분을 어떻게 바꿀지 확정할 수 없다."],
        ["이미지 첨부", "이미지만 보내고 최종 수정 의도를 생략", "현재 상태는 보여주지만 삭제·이동·기능 추가 중 무엇을 원하는지 알 수 없었다.", "산부인과 용어집 115·117·119·121행, 건강증상도우미 63행", "이미지 전송 후 추측 대상과 동작을 확인하는 단계가 필요하다."],
        ["이미지 첨부", "여러 이미지와 번호의 대응 관계가 불완전", "A/B와 색상 번호를 교차 참조하며 이동·삭제·추가를 동시에 지시했다.", "맞춤레시피 55~59행", "참조 그래프가 복잡해져 일부 탭과 검색창이 누락됐다."],
        ["이미지 첨부", "좌우·행·열 공간 표현이 모호", "좌우 배치를 요청했지만 상하로 구현되거나 2열/1열이 다르게 해석됐다.", "인바디코치 75~76행, 금융브리핑 46행", "레이아웃 축과 최종 순서를 함께 적어야 한다."],
        ["이미지 첨부", "부모 컨테이너와 내부 요소의 삭제 경계가 불분명", "박스만 지울지, 내부 기능도 지울지 불분명했다.", "맞춤레시피 57~59행, 출석체크90 135~136행", "삭제 대상과 보존 대상을 쌍으로 명시해야 한다."],
        ["이미지 첨부", "시스템 권한·서비스·위젯 문제를 화면 캡처만으로 해결 시도", "알림 접근 허용 화면과 위젯 오류는 보여줬지만 실제 서비스 상태·매니페스트·런타임 로그는 없었다.", "대화일정 140·141·145·147·151행", "이미지는 '무엇이 보이는지'는 알려주지만 '왜 시스템이 작동하지 않는지'는 알려주지 못한다."],
        ["텍스트", "한 요청에 변경을 과도하게 묶음", "여러 화면·데이터·동작을 한 문단에서 바꾸며 우선순위와 보존 대상을 생략했다.", "건강코치 23행, 금융브리핑 43~47행, 영양분석 109행", "일부 요구만 반영되어 성공 여부를 한 개의 라벨로 판정하기 어려워졌다."],
        ["텍스트", "추상적 동사와 개발 용어를 혼용", "'기능을 합쳐줘', '연동해줘', '2-column'이 사용자가 기대한 화면과 다르게 해석됐다.", "포토스팟 15~16행, 금융브리핑 46행", "완성 상태의 배치·동작으로 바꿔 확인해야 한다."],
        ["텍스트", "과거 상태를 '아까', '기존처럼'으로만 참조", "현재 리비전에 없는 요소나 어느 버전의 상태인지 확정할 수 없었다.", "금융브리핑 43~47행", "리비전 ID와 되돌릴 요소를 함께 전달해야 한다."],
        ["텍스트", "보존 조건·우선순위가 없음", "메인 화면에 남길 요소와 다른 화면으로 옮길 요소가 불분명했다.", "금융브리핑 43~47행, 영양분석 109행", "수정으로 지켜야 할 요소를 명시하지 않으면 회귀가 발생한다."],
        ["텍스트", "'연동'의 데이터 주체·방향·저장 위치가 불분명", "외부 앱을 여는 것과 데이터를 현재 앱 안에 복사하는 것이 구분되지 않았다.", "감정진단 33~36행", "데이터 흐름을 주체→이동→저장→표시로 표현해야 한다."],
        ["텍스트", "이전 실패와 같은 내용을 '다시'로만 재요청", "원인, 변경할 구현, 검증 기준이 추가되지 않았다.", "영양분석 108·110행, 대화일정 153행", "이전 실패를 재현 조건·기대 결과·새로운 제약으로 변환해야 한다."],
        ["공통", "표현은 구체적이지만 빌드·런타임 문제로 실패", "상태 갱신·재실행·수집 기준을 적은 요청이 Flutter build 단계에서 실패했다.", "대화일정 150행", "프롬프트 실패로 분류하면 사용자 표현 특성과 성공률을 왜곡한다."],
        ["공통", "기능 범위가 크지만 최소 성공 범위·폴백이 없음", "위젯에 오늘·내일 일정을 모두 보이려는 요청이 반복 실패했다.", "대화일정 145·147·151행", "핵심 목표를 유지한 최소 범위와 실패 시 줄일 항목을 사전에 정의해야 한다."],
        ["공통", "방식별 난이도 차이를 통제하지 않은 단순 비교", rate_text, "전체 앱 생성·수정 라벨 138건", "이미지 요청은 이미 실패한 복잡 문제에 추가로 사용되는 경향이 있어 단순 비율 차이는 매체 효과가 아닐 수 있다."],
    ]

    style_analysis_sheet(workbook["성공 방식 특징"], success_rows, [15, 30, 54, 39, 52])
    style_analysis_sheet(workbook["실패 방식 특징"], failure_rows, [15, 31, 54, 39, 52])


def update_implications(workbook, app_counts: dict[str, Counter]) -> None:
    image = app_counts[IMAGE_MODE]
    text = app_counts[TEXT_MODE]
    image_n = sum(image.values())
    text_n = sum(text.values())
    rate_summary = (
        f"엄격 성공: 이미지 {image['성공']}/{image_n}({image['성공']/image_n:.1%}), "
        f"텍스트 {text['성공']}/{text_n}({text['성공']/text_n:.1%}); "
        f"일부 실패 포함: 이미지 {(image['성공']+image['일부 실패'])/image_n:.1%}, "
        f"텍스트 {(text['성공']+text['일부 실패'])/text_n:.1%}."
    )
    rows = [
        ["높음", "컴포넌트 선택형 UI 수정 모드", "생성 앱에 안정적인 컴포넌트 ID와 Flutter Semantics/위젯 메타데이터를 심고, 호스트 오버레이에서 터치한 UI를 하이라이트·태깅한다. 정보가 없는 화면은 스크린샷 좌표로 폴백한다.", "표시된 대상+동작 조합은 반복 성공했고, '이렇게'처럼 대상이 모호한 요청은 실패했다. Flutter UI는 일반적으로 Android XML이 아니므로 XML 추출보다 런타임 위젯 정보가 적합하다.", f"대화일정 144행 성공, 감정진단 28~32행 성공. {rate_summary}"],
        ["높음", "표준 심볼 기반 스크린샷 편집기", "사각형=대상, 화살표=이동, X=삭제, 색상=변경, 번호=요구사항 매핑으로 규칙을 고정한다. 전송 전 '대상 ID·동작·목표 위치'로 구조화해 사용자에게 확인시킨다.", "색 동그라미와 화살표를 사용하고 최종 위치를 적은 달력 수정은 성공했다. 반면 복잡한 A/B·색상 교차 참조는 누락을 만들었다.", "대화일정 144행, 맞춤레시피 55~59행, 인바디코치 76행"],
        ["높음", "대상·동작·기대 결과·보존 항목 요청 컴파일러", "자유 입력을 네 항목으로 자동 변환하고, 빈 항목만 짧게 묻는다. 사용자가 확인한 구조화 요청을 Codex와 DB에 동일하게 저장한다.", "성공 요청은 완성 상태와 지켜야 할 요소를 포함했다. '다시', '합쳐줘'는 검증 가능한 결과가 없었다.", "출석체크90 135~136행, 맞춤레시피 57~59행, 대화일정 153행"],
        ["높음", "이미지만 전송하면 의도 확인 후 수정", "AI가 추정한 대상·가능한 동작·보존 요소를 1~3개로 제시하고 사용자 확인 후만 수정을 시작한다.", "이미지만으로는 삭제·이동·오류 설명 중 어느 의도인지 확정할 수 없었다. 일부 '성공'는 질문 생성 성공으로, 앱 수정 성공과 다르다.", "산부인과 용어집 115·117·119·121행, 건강증상도우미 63행"],
        ["높음", "UI 수정과 시스템 기능 문제를 자동 분기", "UI 문제는 시각 태깅을 사용하고, 권한·백그라운드 서비스·위젯·외부 연동은 매니페스트·서비스 상태·로그·재현 절차를 자동 수집한다.", "이미지는 보이는 UI 위치에 강했지만 알림 권한 이미지 2회는 모두 실패했다. 상태 전이 조건을 번호로 정의한 146행은 성공했다.", "대화일정 140·141행 0/2, 146행 1/1. 150행은 구체적이었으나 빌드 오류."],
        ["높음", "복수 변경을 원자적 작업과 항목별 성공 기준으로 분할", "여러 변경이면 번호별 작업으로 나누고, 각 항목에 완료·실패·보류 상태를 저장한다. 일부 실패 시 실패한 항목만 재시도한다.", "현재 '일부 실패'는 어떤 하위 요구가 성공했는지 집계하기 어렵고, 한 문단의 복수 요청에서 누락이 반복됐다.", "건강코치 23행, 금융브리핑 43~47행, 영양분석 109행"],
        ["중간", "실패 시 기능 범위 축소·폴백 제안", "최소 성공 기능을 먼저 만들고 확장하도록 2~3개 범위를 제안한다. 사용자가 허용한 폴백은 요청과 함께 저장한다.", "위젯은 3회 실패한 뒤 표시 범위를 오늘 일정만으로 줄이고 규격을 축소한 요청에서 성공했다.", "대화일정 위젯 시도 145·147·151·154행: 1/4 성공"],
        ["중간", "런타임 오류 증거 번들 자동 첨부", "오류 시점의 현재 화면, 상호작 순서, logcat, 권한/서비스 상태, 앱 버전을 개인정보 제거 후 수정 요청에 자동 첨부한다.", "오류 화면+재현 절차는 UI 오류에서 성공했지만, 시스템 문제는 화면만으로 원인을 알 수 없었다.", "지도 4~6행, 대화일정 140·141·145·147·151행"],
        ["중간", "리비전·보존 대상·컴포넌트 ID를 요청에 자동 첨부", "'아까', '기존처럼'을 사용하면 기준 리비전과 해당 컴포넌트를 선택하게 하고, 보존할 기능을 함께 저장한다.", "상대적 과거 참조는 현재 코드 상태와 다르게 해석되어 회귀를 만들었다.", "금융브리핑 43~47행"],
        ["중간", "성공률을 문제 유형·난이도·기술 실패와 분리해 측정", "UI 이동/삭제, 상태 로직, 외부 연동, 위젯 등으로 층화하고, 빌드·설치·런타임 오류는 '표현 실패'와 별도로 집계한다.", "이미지 엄격 성공률이 텍스트보다 낮지만, 이미지가 복잡하거나 이미 실패한 문제에 사용된 선택 편향이 있다.", rate_summary],
        ["연구", "이미지만 첨부한 사례의 성공을 이중 정의", "'적절한 확인 질문 생성'과 '최종 앱 수정 성공'을 별도 컬럼으로 라벨링한다.", "질문을 잘 한 것은 상호작용 성공이지만 앱 수정 성공은 아니므로 하나의 라벨로 합치면 매체 효과가 과대 평가된다.", "산부인과 용어집 115·117·119·121행"],
    ]
    style_analysis_sheet(workbook["Implication"], rows, [12, 34, 65, 65, 48])


def calculate_app_counts(rows: list[dict[str, object]]) -> dict[str, Counter]:
    result = {IMAGE_MODE: Counter(), TEXT_MODE: Counter()}
    for item in rows:
        if clean(item.get("상호작용 단계")) != APP_STAGE:
            continue
        mode = clean(item.get("분석 방식"))
        label = clean(item.get("성공 판정"))
        if mode in result and label in LABELS:
            result[mode][label] += 1
    return result


def update_workbook() -> dict[str, Counter]:
    if not BACKUP_PATH.exists():
        shutil.copy2(WORKBOOK_PATH, BACKUP_PATH)
    workbook = load_workbook(WORKBOOK_PATH)
    rows = interaction_rows(workbook)
    app_counts = calculate_app_counts(rows)
    update_task_collection(workbook, rows)
    update_characteristic_sheets(workbook, app_counts)
    update_implications(workbook, app_counts)
    workbook.save(WORKBOOK_PATH)
    return app_counts


BLACK = RGBColor(0x16, 0x1B, 0x1D)
MUTED = RGBColor(0x5E, 0x68, 0x6A)
TEAL = RGBColor(0x15, 0x81, 0x58)
BLUE = RGBColor(0x05, 0x8D, 0xC7)
ORANGE = RGBColor(0xED, 0x56, 0x1B)
PALE = RGBColor(0xF3, 0xF3, 0xF3)
LIGHT_TEAL = RGBColor(0xE8, 0xF3, 0xEF)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xF9)
LIGHT_ORANGE = RGBColor(0xFC, 0xEE, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"


def remove_all_slides(prs: Presentation) -> None:
    for slide_id in list(prs.slides._sldIdLst):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def add_text(slide, x: float, y: float, w: float, h: float, text: str, *, size: float = 18,
             color: RGBColor = BLACK, bold: bool = False, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, margin: float = 0.04) -> object:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_bullets(slide, x: float, y: float, w: float, h: float, bullets: list[str], *, size: float = 16,
                color: RGBColor = BLACK, accent: RGBColor = TEAL) -> object:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    for index, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(7)
        paragraph.level = 0
    return shape


def add_rect(slide, x: float, y: float, w: float, h: float, fill: RGBColor, *, line: RGBColor | None = None,
             radius: bool = False) -> object:
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if line is None:
        shape.line.fill.background()
    return shape


def add_title(slide, title: str, number: int, subtitle: str | None = None) -> None:
    add_text(slide, 0.34, 0.37, 8.9, 0.55, title, size=25, bold=True)
    if subtitle:
        add_text(slide, 0.35, 0.91, 8.8, 0.26, subtitle, size=10.5, color=MUTED)
    add_text(slide, 9.28, 5.08, 0.45, 0.25, str(number), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_new_slide(prs: Presentation, title: str, number: int, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[10])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_title(slide, title, number, subtitle)
    return slide


def add_picture_contain(slide, image_path: Path, x: float, y: float, w: float, h: float) -> object:
    with Image.open(image_path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    display_w = iw * scale
    display_h = ih * scale
    return slide.shapes.add_picture(
        str(image_path),
        Inches(x + (w - display_w) / 2),
        Inches(y + (h - display_h) / 2),
        Inches(display_w),
        Inches(display_h),
    )


def add_rate_bar(slide, y: float, name: str, counts: Counter, accent: RGBColor) -> None:
    total = sum(counts.values())
    success = counts["성공"] / total
    partial = counts["일부 실패"] / total
    failure = counts["실패"] / total
    add_text(slide, 0.55, y - 0.04, 1.2, 0.32, name, size=17, bold=True)
    add_text(slide, 1.65, y - 0.08, 1.25, 0.42, f"{success:.1%}", size=26, bold=True, color=accent)
    x = 2.85
    width = 6.35
    add_rect(slide, x, y, width * success, 0.36, accent)
    add_rect(slide, x + width * success, y, width * partial, 0.36, RGBColor(0xF0, 0xB4, 0x29))
    add_rect(slide, x + width * (success + partial), y, width * failure, 0.36, ORANGE)
    add_text(slide, 2.85, y + 0.39, 6.3, 0.25,
             f"성공 {counts['성공']}/{total}  ·  일부 실패 {counts['일부 실패']}/{total}  ·  실패 {counts['실패']}/{total}",
             size=10.5, color=MUTED)


def build_presentation(app_counts: dict[str, Counter]) -> None:
    reference_pptx = find_reference_pptx()
    prs = Presentation(reference_pptx)
    remove_all_slides(prs)

    image = app_counts[IMAGE_MODE]
    text = app_counts[TEXT_MODE]
    total = sum(image.values()) + sum(text.values())

    slide = add_new_slide(prs, "이미지 첨부 방식과 텍스트 방식의 수정 성공 분석", 1)
    add_text(slide, 0.55, 1.45, 8.75, 1.1, "참가자가 어떻게 표현했을 때\n생성 앱의 수정이 성공했는가", size=27, bold=True)
    add_text(slide, 0.58, 3.05, 8.2, 0.55, f"앱 생성·수정 요청 라벨 {total}건 분석", size=18, color=TEAL, bold=True)
    add_text(slide, 0.58, 3.68, 8.25, 0.85, "성공 / 일부 실패 / 실패 라벨과 요청 원문·첨부 이미지를 함께 검토\n2026. 08. 10.", size=14, color=MUTED)

    slide = add_new_slide(prs, "방식별 성공 비율", 2, "엄격 성공률 = '성공'만 포함; 일부 실패는 별도 표시")
    add_rate_bar(slide, 1.55, "이미지", image, TEAL)
    add_rate_bar(slide, 2.75, "텍스트", text, BLUE)
    add_rect(slide, 0.55, 3.9, 8.7, 0.72, PALE)
    add_text(slide, 0.75, 4.04, 8.25, 0.4,
             "텍스트 성공률이 8.7%p 높지만, 이미지는 복잡하거나 이미 실패한 문제에 추가로 사용된 경우가 있어 인과효과로 단정할 수 없다.",
             size=13.5)

    slide = add_new_slide(prs, "이미지가 성공한 방식", 3, "이미지는 '어디', 텍스트는 '무엇을·어떻게'를 정확하게 했을 때 보완적으로 작동")
    add_rect(slide, 0.45, 1.25, 3.0, 3.65, LIGHT_TEAL)
    add_picture_contain(slide, IMAGES_DIR / "068.jpg", 0.65, 1.42, 2.6, 2.55)
    add_text(slide, 0.67, 4.05, 2.55, 0.65, "대화일정 144행 · 성공\n색 표시 + 최종 배치 명시", size=12.5, bold=True)
    add_bullets(slide, 3.78, 1.35, 5.7, 2.3, [
        "화면 대상을 색·동그라미·번호로 특정",
        "삭제·이동·색상 변경처럼 한정된 동사 사용",
        "좌우·상하·행·열과 요소 순서로 완성 상태 설명",
        "복수 이미지는 번호와 요구사항을 1:1로 매핑",
    ], size=15.2)
    add_rect(slide, 3.85, 4.0, 5.2, 0.68, PALE)
    add_text(slide, 4.02, 4.12, 4.9, 0.4, "같은 달력 UI는 구체적 텍스트만 사용한 149행에서도 성공", size=12.5)

    slide = add_new_slide(prs, "이미지가 실패한 방식", 4, "보이는 현상과 시스템 내부 원인은 다른 정보")
    add_picture_contain(slide, IMAGES_DIR / "067.jpg", 0.5, 1.22, 2.35, 2.25)
    add_picture_contain(slide, IMAGES_DIR / "071.jpg", 2.95, 1.22, 2.35, 2.25)
    add_text(slide, 0.58, 3.52, 2.2, 0.54, "알림 권한 이미지\n2회 모두 실패", size=13, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 3.03, 3.52, 2.2, 0.54, "위젯 연속 시도\n4회 중 1회 성공", size=13, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, 5.6, 1.28, 3.8, 2.65, [
        "'이렇게', '다시'처럼 의도를 이미지에 맡김",
        "위젯·권한·서비스를 스크린샷만으로 진단",
        "복수 대상·동작·이미지 번호가 교차 참조",
        "부모 박스와 내부 요소의 삭제/보존 경계 생략",
    ], size=14.5)
    add_rect(slide, 5.62, 4.05, 3.72, 0.62, LIGHT_ORANGE)
    add_text(slide, 5.78, 4.16, 3.38, 0.36, "이미지는 '무엇이 보이는지'는 알려주지만 '왜'는 보장하지 않음", size=11.5)

    slide = add_new_slide(prs, "Implication 1. 터치로 수정 대상을 선택", 5, "사용자가 코드 변수명을 모르더라도 '이 UI'를 정확하게 지정")
    add_rect(slide, 0.55, 1.35, 2.4, 3.25, PALE, line=MUTED, radius=True)
    add_text(slide, 0.8, 1.65, 1.9, 0.4, "생성 앱 화면", size=17, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.88, 2.25, 1.75, 0.55, LIGHT_BLUE, line=BLUE, radius=True)
    add_text(slide, 1.03, 2.38, 1.45, 0.25, "수정할 UI", size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.88, 3.05, 1.75, 0.45, WHITE, line=RGBColor(0xB7, 0xBE, 0xC0), radius=True)
    add_rect(slide, 0.88, 3.72, 1.75, 0.45, WHITE, line=RGBColor(0xB7, 0xBE, 0xC0), radius=True)
    add_text(slide, 3.3, 1.42, 5.9, 0.45, "XML 추출보다 Flutter 런타임 정보가 적합", size=19, color=TEAL, bold=True)
    add_bullets(slide, 3.32, 2.0, 5.9, 2.5, [
        "생성 시 안정적인 component_id와 Semantics/위젯 메타데이터를 삽입",
        "호스트 오버레이에서 터치한 UI를 하이라이트·태깅",
        "Codex에 화면 좌표가 아닌 컴포넌트 ID, 현재 속성, 부모 관계를 전달",
        "메타데이터가 없는 UI는 스크린샷 좌표로 폴백",
    ], size=14.4)

    slide = add_new_slide(prs, "Implication 2. 심볼의 의미를 고정한 스크린샷 편집기", 6, "그림을 추측하지 않고 구조화된 요청으로 변환")
    symbols = [
        ("□", "대상 지정", TEAL),
        ("→", "이동", BLUE),
        ("X", "삭제", ORANGE),
        ("1", "요구사항 매핑", RGBColor(0x7A, 0x56, 0xA6)),
    ]
    for index, (symbol, label, color) in enumerate(symbols):
        x = 0.55 + index * 2.28
        add_rect(slide, x, 1.35, 1.95, 1.0, WHITE, line=color, radius=True)
        add_text(slide, x + 0.12, 1.5, 0.48, 0.48, symbol, size=25, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.62, 1.62, 1.17, 0.28, label, size=13, bold=True)
    add_rect(slide, 0.55, 2.72, 8.85, 1.23, PALE)
    add_text(slide, 0.78, 2.91, 8.35, 0.3, "전송 전 변환 결과", size=14, color=TEAL, bold=True)
    add_text(slide, 0.78, 3.24, 8.3, 0.45, "대상: calendar.header.month  |  동작: 이동  |  목표: 상단 중앙  |  보존: 왼·오른쪽 화살표", size=14)
    add_text(slide, 0.65, 4.25, 8.6, 0.45, "근거: 대상·동작·최종 배치가 함께 제공된 144행은 성공. 복잡한 A/B·색상 교차 참조는 일부 누락.", size=13, color=MUTED)

    slide = add_new_slide(prs, "Implication 3. 자유 입력을 구조화한 수정 요청으로", 7, "사용자 문장을 강제로 바꾸지 않고, 빈 정보만 확인")
    labels = [
        ("대상", "어느 화면·UI·기능인가", LIGHT_TEAL, TEAL),
        ("동작", "삭제·이동·추가·수정", LIGHT_BLUE, BLUE),
        ("기대 결과", "완료 후 무엇이 보이고 작동하나", RGBColor(0xF7, 0xF2, 0xE2), RGBColor(0x9C, 0x71, 0x16)),
        ("보존", "바꾸지 말아야 할 요소는 무엇인가", LIGHT_ORANGE, ORANGE),
    ]
    for index, (heading, body, fill, color) in enumerate(labels):
        x = 0.5 + index * 2.35
        add_rect(slide, x, 1.4, 2.08, 1.65, fill)
        add_text(slide, x + 0.16, 1.6, 1.75, 0.35, heading, size=17, bold=True, color=color)
        add_text(slide, x + 0.16, 2.08, 1.75, 0.68, body, size=12.4)
    add_text(slide, 0.55, 3.42, 4.15, 0.4, "복수 변경은 번호별 작업으로 분할", size=18, bold=True)
    add_bullets(slide, 0.58, 3.88, 4.2, 0.82, ["항목별 성공/실패 저장", "실패한 항목만 재시도"], size=13.3)
    add_rect(slide, 5.0, 3.4, 4.35, 1.15, PALE)
    add_text(slide, 5.2, 3.58, 3.98, 0.72, "'기능을 합쳐줘'·'다시 만들어줘'는 실패\n완성 상태·수용 기준을 나열한 요청은 성공", size=13.2)

    slide = add_new_slide(prs, "Implication 4. 문제 유형에 맞는 정보를 자동 수집", 8, "스크린샷이 답할 수 있는 문제와 시스템 증거가 필요한 문제를 분리")
    add_rect(slide, 0.55, 1.35, 3.95, 2.25, LIGHT_TEAL)
    add_text(slide, 0.8, 1.6, 3.45, 0.35, "UI 배치·색·삭제", size=19, color=TEAL, bold=True)
    add_bullets(slide, 0.82, 2.05, 3.35, 1.25, ["컴포넌트 선택/하이라이트", "스크린샷 심볼", "최종 배치 확인"], size=13.5)
    add_rect(slide, 4.85, 1.35, 4.55, 2.25, LIGHT_ORANGE)
    add_text(slide, 5.1, 1.6, 4.05, 0.35, "권한·서비스·외부 연동·위젯", size=19, color=ORANGE, bold=True)
    add_bullets(slide, 5.12, 2.05, 3.95, 1.25, ["매니페스트·서비스/권한 상태", "logcat·재현 절차·앱 버전", "데이터 주체→저장→표시 흐름"], size=13.5)
    add_text(slide, 0.65, 3.98, 8.6, 0.64, "대화일정: 권한 스크린샷 2/2 실패 → 상태 갱신·수집·복귀 조건을 정의한 146행 1/1 성공.\n동일한 요청이 150행에서는 빌드 오류로 실패했으므로 프롬프트 실패와 기술 실패를 분리해야 한다.", size=13.5)

    slide = add_new_slide(prs, "추가 Implication", 9, "실패을 다시 같이 반복하지 않도록 작업 범위와 증거를 관리")
    items = [
        ("범위 축소·폴백", "실패 시 최소 기능을 제안하고 성공 후 확장", "위젯 4회 중 범위를 줄인 1회만 성공"),
        ("리비전·보존 자동 첨부", "'아까'를 현재 리비전·컴포넌트 ID로 변환", "과거 상대 참조에서 회귀·누락 발생"),
        ("이미지만 전송 확인", "추정한 대상·동작을 확인한 후 수정 시작", "질문 성공과 최종 수정 성공을 분리"),
        ("항목별 성공 라벨", "일부 실패를 하위 요구별로 저장", "누락된 항목만 재시도·분석 가능"),
    ]
    for index, (heading, body, evidence) in enumerate(items):
        row, col = divmod(index, 2)
        x = 0.55 + col * 4.48
        y = 1.32 + row * 1.75
        add_rect(slide, x, y, 4.15, 1.42, PALE)
        add_text(slide, x + 0.18, y + 0.16, 3.78, 0.28, heading, size=16, color=TEAL if col == 0 else BLUE, bold=True)
        add_text(slide, x + 0.18, y + 0.53, 3.78, 0.34, body, size=12.3)
        add_text(slide, x + 0.18, y + 1.0, 3.78, 0.27, evidence, size=10.6, color=MUTED)

    slide = add_new_slide(prs, "결론과 구현 우선순위", 10)
    add_text(slide, 0.6, 1.2, 8.8, 0.55, "이미지 첨부 자체보다 의도를 구조화하는 상호작용이 핵심", size=23, bold=True)
    priorities = [
        ("1", "컴포넌트 선택 오버레이 + 심볼 편집기", "UI의 '어디'를 모호하지 않게 전달"),
        ("2", "대상·동작·기대 결과·보존 요청 컴파일러", "자유 입력을 검증 가능한 요청으로 변환"),
        ("3", "UI/시스템 문제 분기 + 런타임 증거 번들", "권한·위젯·외부 연동 문제의 진단 근거 확보"),
        ("4", "항목별 성공·기술 실패 분리 측정", "기능의 효과를 왜곡 없이 검증"),
    ]
    for index, (number, heading, body) in enumerate(priorities):
        y = 1.95 + index * 0.76
        add_rect(slide, 0.7, y, 0.48, 0.48, TEAL if index < 2 else BLUE, radius=True)
        add_text(slide, 0.71, y + 0.08, 0.46, 0.25, number, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, 1.4, y - 0.02, 4.85, 0.31, heading, size=16, bold=True)
        add_text(slide, 6.22, y - 0.01, 3.03, 0.45, body, size=11.8, color=MUTED)

    prs.save(OUTPUT_PPTX)


def validate(app_counts: dict[str, Counter]) -> None:
    workbook = load_workbook(WORKBOOK_PATH, read_only=True, data_only=False)
    assert workbook["Task별 방식 모음"].max_row >= 22
    assert workbook["성공 방식 특징"].max_row >= 10
    assert workbook["실패 방식 특징"].max_row >= 10
    assert workbook["Implication"].max_row >= 10
    classified = workbook["상호작용 분류"]
    assert classified.cell(140, 9).value == "실패"
    assert classified.cell(154, 9).value == "성공"
    presentation = Presentation(OUTPUT_PPTX)
    assert len(presentation.slides) == 10
    assert presentation.slide_width == Inches(10)
    assert presentation.slide_height == Inches(5.625)
    assert sum(app_counts[IMAGE_MODE].values()) == 64
    assert sum(app_counts[TEXT_MODE].values()) == 74


def main() -> None:
    app_counts = update_workbook()
    build_presentation(app_counts)
    validate(app_counts)
    print(f"Updated workbook: {WORKBOOK_PATH}")
    print(f"Backup: {BACKUP_PATH}")
    print(f"Created deck: {OUTPUT_PPTX}")
    for mode, counts in app_counts.items():
        print(mode, dict(counts), "total", sum(counts.values()))


if __name__ == "__main__":
    main()
