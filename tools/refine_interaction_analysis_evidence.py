#!/usr/bin/env python3
"""Add participant text and quantitative evidence to the analysis workbook/deck."""

from __future__ import annotations

from collections import Counter
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from update_interaction_analysis_artifacts import (
    BLACK,
    BLUE,
    EXPORTS,
    FONT,
    IMAGES_DIR,
    LIGHT_BLUE,
    LIGHT_ORANGE,
    LIGHT_TEAL,
    MUTED,
    ORANGE,
    OUTPUT_PPTX,
    PALE,
    TEAL,
    WHITE,
    WORKBOOK_PATH,
    add_bullets,
    add_new_slide,
    add_picture_contain,
    add_rate_bar,
    add_rect,
    add_text,
    find_reference_pptx,
    remove_all_slides,
)


LABELS = ("성공", "일부 실패", "실패")


def clean(value: object) -> str:
    return str(value or "").strip()


def load_labeled_rows(workbook) -> dict[int, dict[str, str]]:
    ws = workbook["상호작용 분류"]
    result: dict[int, dict[str, str]] = {}
    for row in range(2, ws.max_row + 1):
        label = clean(ws.cell(row, 9).value)
        source_row = ws.cell(row, 13).value
        if label not in LABELS or source_row in (None, ""):
            continue
        result[int(source_row)] = {
            "task": clean(ws.cell(row, 2).value),
            "stage": clean(ws.cell(row, 6).value),
            "mode": clean(ws.cell(row, 8).value),
            "label": label,
            "text": clean(ws.cell(row, 10).value),
            "evidence": clean(ws.cell(row, 11).value),
            "files": clean(ws.cell(row, 12).value),
        }
    return result


def count_rows(rows: dict[int, dict[str, str]], source_rows: list[int]) -> Counter:
    return Counter(rows[row]["label"] for row in source_rows if row in rows)


def metric(counts: Counter) -> tuple[int, int, float, float]:
    total = sum(counts.values())
    success = counts["성공"]
    strict = success / total if total else 0.0
    nonfailure = (success + counts["일부 실패"]) / total if total else 0.0
    return success, total, strict, nonfailure


def append_once(text: str, addition: str) -> str:
    if addition in text:
        return text
    return f"{text}\n\n{addition}" if text else addition


def update_workbook_evidence(workbook, rows: dict[int, dict[str, str]]) -> None:
    success_ws = workbook["성공 방식 특징"]
    success_example = (
        '참가자 원문 예시(144행): "빨간 동그라미에 있는 년월을 달력 상단 중앙에 배치해주고… '
        '파란 동그라미에 있는 왼쪽 화살표를… 년월 왼쪽에… 노란색 동그라미에 있는 오른쪽 화살표도… 년월 오른쪽에 배치해줘."'
    )
    for row in range(2, success_ws.max_row + 1):
        characteristic = clean(success_ws.cell(row, 2).value)
        if characteristic == "표시된 대상과 수정 동작을 함께 명시":
            success_ws.cell(row, 3).value = append_once(clean(success_ws.cell(row, 3).value), success_example)
            success_ws.cell(row, 5).value = append_once(
                clean(success_ws.cell(row, 5).value),
                "선별한 '시각 표식+명시적 UI 동작' 14건 중 성공 10건(71.4%), 일부 실패 포함 13건(92.9%). 원본 7·10·41~45·47·52·68·74~76·144행.",
            )
        elif characteristic == "이미지·표시를 번호와 요구사항으로 1:1 대응":
            success_ws.cell(row, 3).value = append_once(
                clean(success_ws.cell(row, 3).value),
                '참가자 원문 예시(76행): "#1 이미지_붉은색 동그라미 안에 있는 아이콘을 2열이 아니라 1열 좌우로 배치… #2 단기목표·장기목표… #3 운동가능시간…"',
            )
        elif characteristic == "삭제 대상과 보존 대상의 경계를 명시":
            success_ws.cell(row, 3).value = append_once(
                clean(success_ws.cell(row, 3).value),
                '참가자 원문 예시(136행): "스크린샷처럼 시작일 설정과 날짜별 기록 수정 두 버튼만 남기고 출석체크 90의 큰 박스와 내용물은 제거해줘."',
            )

    failure_ws = workbook["실패 방식 특징"]
    failure_example = (
        '참가자 원문 예시(141행): "진한 빨간 동그라미 안에 있는 버튼을 눌렀을 때… '
        '권한 허용을 하여도 알림 접근 화면에 변화가 일어나질 않아… 체크 표시를 해줬으면 좋겠어."'
    )
    for row in range(2, failure_ws.max_row + 1):
        characteristic = clean(failure_ws.cell(row, 2).value)
        if characteristic == "'이렇게', '다시'처럼 이미지 해석에 전적으로 의존":
            failure_ws.cell(row, 3).value = append_once(
                clean(failure_ws.cell(row, 3).value),
                '참가자 원문 예시: 108행 "이렇게 수정. 사진 이해 못했으면 말해" / 110행 "다시. 화나게 하지 말고". 두 요청 모두 실패.',
            )
        elif characteristic == "부모 컨테이너와 내부 요소의 삭제 경계가 불분명":
            failure_ws.cell(row, 3).value = append_once(
                clean(failure_ws.cell(row, 3).value),
                '참가자 원문 예시(135행): "출석체크 90 박스 안의 내용을 스크린샷처럼 변경해줘." → 박스·내용·버튼 중 삭제/보존 경계가 모호해 요청하지 않은 텍스트가 추가된 일부 실패.',
            )
        elif characteristic == "시스템 권한·서비스·위젯 문제를 화면 캡처만으로 해결 시도":
            failure_ws.cell(row, 3).value = append_once(clean(failure_ws.cell(row, 3).value), failure_example)
            failure_ws.cell(row, 5).value = append_once(
                clean(failure_ws.cell(row, 5).value),
                "시스템 상태를 이미지로 보여준 선별 사례 5건(140·141·145·147·151행)은 성공 0건(0%). 화면 밖 로그·권한·서비스 상태가 필요했다.",
            )

    implication_metrics = {
        "컴포넌트 선택형 UI 수정 모드": (
            "참가자들은 수정할 UI 위에 원·색·번호·블록을 그리고 '삭제', '좌우 배치', '상단 중앙으로 이동'같은 텍스트를 덧붙였다. 이 조합은 선별 14건 중 10건 성공(71.4%), 일부 실패 포함 13건(92.9%)이었다.",
            '원본 7·10·41~45·47·52·68·74~76·144행. 예: 144행 "빨간 동그라미의 년월을 달력 상단 중앙에…"',
        ),
        "표준 심볼 기반 스크린샷 편집기": (
            "참가자들은 1·2·3번, 색상, A/B를 사용했지만 의미가 요청마다 달랐다. 임의 번호 매핑 9건의 엄격 성공은 4건(44.4%), 일부 실패 포함 7건(77.8%)이어서 심볼 의미를 표준화할 필요가 있다.",
            '원본 42~45·47·58·68·76·83행. 예: 42행 "2번 블록을 메인화면에서 지워… 1번 아이콘을 누르면 따로 띄우게…"',
        ),
        "대상·동작·기대 결과·보존 항목 요청 컴파일러": (
            "최종 배치와 보존 대상을 구체적으로 적은 UI 사례는 16건 중 14건 성공(87.5%)이었다. 반면 '이렇게', '수정해줘'처럼 의도가 모호한 이미지 요청 6건은 엄격 성공 0건(0%)이었다.",
            '구체 배치 24·29·41·42·45·46·59·65·67·68·75·76·85·134·136·144행. 모호 요청 58·63·82·108·110·135행.',
        ),
        "이미지만 전송하면 의도 확인 후 수정": (
            "이미지만 전송한 5건은 현재 라벨에서 성공 4건(80%), 일부 실패 1건이었다. 다만 여기서 '성공'은 여러 경우 앱 수정이 아니라 AI가 적절한 확인 질문을 생성했다는 뜻이었다.",
            "원본 95·115·117·119·121행. 입력 원문은 모두 '(이미지만 첨부)'. 질문 성공과 수정 성공을 분리해야 함.",
        ),
        "UI 수정과 시스템 기능 문제를 자동 분기": (
            "권한·위젯 시스템 상태를 이미지로만 설명한 선별 5건은 성공 0건(0%)이었다. 반면 권한 갱신·서비스 연결·수집 조건을 번호로 정의한 146행은 성공했다.",
            '시스템 이미지 140·141·145·147·151행 0/5. 예: 141행 "권한 허용을 하여도… 변화가 일어나지 않아." 146행은 9개 수용 기준을 제시.',
        ),
        "복수 변경을 원자적 작업과 항목별 성공 기준으로 분할": (
            "번호를 2개 이상 사용한 복수 요청 18건 중 13건이 성공(72.2%)했지만 5건(27.8%)은 일부 실패 또는 실패였다. 현재의 한 개 라벨만으로는 어느 하위 항목이 누락됐는지 알 수 없다.",
            '원본 7·8·10~12·41~45·71·83·87·89·90·93·146·150행. 예: 43행 "1번 블록은 삭제… 2번 블록은 간소화… 3번 블록의 빨간 표시 삭제" → 일부 실패.',
        ),
        "실패 시 기능 범위 축소·폴백 제안": (
            "대화일정 위젯은 4회 시도 중 성공 1회(25%)였다. 처음 3회는 오늘·내일 일정과 원인 추정을 반복했고, 마지막에 '규격을 줄이고 불가능하면 오늘 일정만'으로 폴백을 허용한 뒤 성공했다.",
            '원본 145·147·151·154행. 예: 154행 "위젯 규격을 줄여서… 정보가 너무 많아 불가능하다면 오늘 일정만 보이게…"',
        ),
        "런타임 오류 증거 번들 자동 첨부": (
            "UI 런타임 오류는 오류 화면+재현 동작을 제공한 6행에서 1/1 성공했다. 반면 권한·위젯 오류 이미지 4건(140·141·147·151)은 0/4 성공이어서 문제 유형별 증거가 달라야 함을 보여준다.",
            '예: 6행 "오류 화면을 캡처해서 보낼게… 다시 오류가 나지 않게 수정해줘." 시스템 문제에는 logcat·권한·서비스 상태가 추가로 필요.',
        ),
        "리비전·보존 대상·컴포넌트 ID를 요청에 자동 첨부": (
            "금융브리핑 43~47행은 '아까', 'v2에서 다시', '메인화면엔 그대로'처럼 과거 상태와 보존 대상을 상대적으로 참조했다. 5건 중 엄격 성공은 1건(20%), 일부 실패 4건이었다.",
            '원본 43~47행. 예: 44행 "금융브리핑 v2에서 다시 시작할게… 아까 삭제했던 부분에… 메인화면엔 냅두고…" → 일부 실패.',
        ),
        "성공률을 문제 유형·난이도·기술 실패와 분리해 측정": (
            "전체 엄격 성공률은 이미지 42/64(65.6%), 텍스트 55/74(74.3%)이다. 하지만 같은 권한 수용 기준이 146행에서는 성공하고 150행에서는 Flutter build 오류로 실패했으므로 매체·표현·기술 실패를 분리해야 한다.",
            "전체 앱 생성·수정 라벨 138건. 146행과 150행은 거의 동일한 권한 요구사항이지만 빌드 결과가 다름.",
        ),
        "이미지만 첨부한 사례의 성공을 이중 정의": (
            "이미지만 전송한 5건의 현재 라벨은 성공 4건(80%), 일부 실패 1건이지만, 다수의 '성공'은 앱 수정 완료가 아니라 AI의 확인 질문 생성을 뜻했다. 이를 최종 수정 성공률에 합치면 이미지 효과가 과대 평가된다.",
            "원본 95·115·117·119·121행. '질문 적절성'과 '앱 반영 성공'을 별도 컬럼으로 라벨링해야 함.",
        ),
    }

    implication_ws = workbook["Implication"]
    for row in range(2, implication_ws.max_row + 1):
        title = clean(implication_ws.cell(row, 2).value)
        if title not in implication_metrics:
            continue
        reason, evidence = implication_metrics[title]
        implication_ws.cell(row, 4).value = reason
        implication_ws.cell(row, 5).value = evidence


def add_metric_block(slide, value: str, label: str, *, color: RGBColor = TEAL) -> None:
    add_rect(slide, 0.55, 1.34, 2.55, 1.05, PALE)
    compact_value = value.replace("\n", " · ")
    value_size = 20 if len(compact_value) <= 14 else 17
    add_text(slide, 0.67, 1.5, 2.3, 0.34, compact_value, size=value_size, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 0.7, 1.94, 2.24, 0.28, label, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)


def add_quote_box(slide, text: str, source: str, *, x: float = 3.45, y: float = 1.35, w: float = 5.95, h: float = 1.45) -> None:
    add_rect(slide, x, y, w, h, PALE)
    add_text(slide, x + 0.18, y + 0.14, w - 0.36, h - 0.46, f'"{text}"', size=12.5)
    add_text(slide, x + 0.18, y + h - 0.28, w - 0.36, 0.18, source, size=9.5, color=MUTED, align=PP_ALIGN.RIGHT)


def add_pattern_example_slide(
    prs: Presentation,
    page: int,
    heading: str,
    characteristic: str,
    images: list[Path],
    participant_text: str,
    source: str,
    result: str,
    interpretation: str,
    *,
    success: bool,
) -> None:
    accent = TEAL if success else ORANGE
    slide = add_new_slide(prs, heading, page, characteristic)
    image_area_x, image_area_y, image_area_w, image_area_h = 0.48, 1.25, 3.05, 3.55
    add_rect(slide, image_area_x, image_area_y, image_area_w, image_area_h, LIGHT_TEAL if success else LIGHT_ORANGE)
    count = len(images)
    if count == 1:
        add_picture_contain(slide, images[0], image_area_x + 0.18, image_area_y + 0.18, image_area_w - 0.36, 2.55)
    else:
        gap = 0.08
        each_w = (image_area_w - 0.36 - gap * (count - 1)) / count
        for index, image_path in enumerate(images):
            add_picture_contain(slide, image_path, image_area_x + 0.18 + index * (each_w + gap), image_area_y + 0.23, each_w, 2.45)
    add_text(slide, image_area_x + 0.18, 4.05, image_area_w - 0.36, 0.48, f"{source}\n{result}", size=11.5, color=accent, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, 3.82, 1.3, 5.35, 0.28, "참가자가 이미지와 함께 전송한 텍스트", size=14, color=accent, bold=True)
    add_rect(slide, 3.72, 1.7, 5.68, 1.75, PALE)
    add_text(slide, 3.95, 1.92, 5.2, 1.3, f'"{participant_text}"', size=12.5)
    add_rect(slide, 3.72, 3.72, 5.68, 0.87, LIGHT_TEAL if success else LIGHT_ORANGE)
    add_text(slide, 3.92, 3.86, 1.12, 0.24, "해석", size=12.5, color=accent, bold=True)
    add_text(slide, 4.86, 3.82, 4.3, 0.52, interpretation, size=12.1)


def add_implication_slide(
    prs: Presentation,
    page: int,
    number: int,
    title: str,
    subtitle: str,
    metric_value: str,
    metric_label: str,
    observed: list[str],
    quote: str,
    quote_source: str,
    application: str,
    *,
    color: RGBColor = TEAL,
    image: Path | None = None,
) -> None:
    slide = add_new_slide(prs, f"Implication {number}. {title}", page, subtitle)
    add_metric_block(slide, metric_value, metric_label, color=color)
    add_text(slide, 0.6, 2.62, 2.35, 0.28, "왜 도출했는가", size=15, color=color, bold=True)
    add_bullets(slide, 0.62, 2.98, 2.52, 1.62, observed, size=11.8)

    if image:
        add_rect(slide, 3.42, 1.32, 1.7, 2.1, WHITE, line=RGBColor(0xD0, 0xD5, 0xD6))
        add_picture_contain(slide, image, 3.54, 1.42, 1.46, 1.82)
        add_quote_box(slide, quote, quote_source, x=5.28, y=1.32, w=4.12, h=2.1)
    else:
        add_quote_box(slide, quote, quote_source, x=3.42, y=1.32, w=5.98, h=2.1)
    add_rect(slide, 3.42, 3.72, 5.98, 0.87, LIGHT_TEAL if color == TEAL else LIGHT_BLUE)
    add_text(slide, 3.62, 3.86, 1.1, 0.24, "설계 적용", size=12.5, color=color, bold=True)
    add_text(slide, 4.66, 3.83, 4.52, 0.5, application, size=12.2)


def build_deck(rows: dict[int, dict[str, str]]) -> None:
    prs = Presentation(find_reference_pptx())
    remove_all_slides(prs)

    image_counts = Counter()
    text_counts = Counter()
    for item in rows.values():
        if item["stage"] != "앱 생성·수정 요청":
            continue
        if item["mode"] == "이미지 첨부":
            image_counts[item["label"]] += 1
        elif item["mode"] == "텍스트":
            text_counts[item["label"]] += 1

    slide = add_new_slide(prs, "이미지 첨부 방식과 텍스트 방식의 수정 성공 분석", 1)
    add_text(slide, 0.55, 1.42, 8.8, 1.05, "참가자의 표현 양상, 수정 성공률,\n상호작용 설계 Implication", size=28, bold=True)
    add_text(slide, 0.58, 3.12, 8.0, 0.38, "앱 생성·수정 요청 라벨 138건", size=18, color=TEAL, bold=True)
    add_text(slide, 0.58, 3.72, 8.3, 0.8, "성공 / 일부 실패 / 실패 라벨·요청 원문·첨부 이미지를 함께 검토\n2026. 08. 10.", size=14, color=MUTED)

    slide = add_new_slide(prs, "방식별 성공 비율", 2, "엄격 성공률 = '성공'만 포함; 일부 실패는 별도 표시")
    add_rate_bar(slide, 1.55, "이미지", image_counts, TEAL)
    add_rate_bar(slide, 2.75, "텍스트", text_counts, BLUE)
    add_rect(slide, 0.55, 3.9, 8.7, 0.72, PALE)
    add_text(slide, 0.75, 4.04, 8.25, 0.4, "텍스트가 8.7%p 높지만, 이미지는 복잡하거나 이미 실패한 문제에 추가로 사용된 경우가 있어 매체의 인과효과로 단정하지 않음.", size=13.2)

    add_pattern_example_slide(
        prs, 3, "이미지가 성공한 방식 1", "표시된 대상과 수정 동작을 함께 명시",
        [IMAGES_DIR / "068.jpg"],
        "빨간 동그라미에 있는 년월을 달력 상단 중앙에 배치해주고, 파란 동그라미에 있는 왼쪽 화살표를 년월 왼쪽에, 노란색 동그라미에 있는 오른쪽 화살표를 년월 오른쪽에 배치해줘.",
        "P01 · 대화일정 144행", "성공",
        "색 표식이 대상을 정하고, 텍스트가 이동 방향·최종 위치·버튼 동작을 확정했다.", success=True,
    )
    add_pattern_example_slide(
        prs, 4, "이미지가 성공한 방식 2", "이미지·표시를 번호와 요구사항으로 1:1 대응",
        [IMAGES_DIR / "035.jpg", IMAGES_DIR / "036.jpg", IMAGES_DIR / "037.jpg"],
        "#1 이미지_붉은색 동그라미 안에 있는 아이콘을 2열이 아닌 1열 좌우로 배치해줘. #2 단기목표·장기목표는 내용을 기입할 수 있게… #3 운동가능시간에 60분·90분·120분 이상을 추가해줘.",
        "P01 · 인바디코치 76행", "성공",
        "이미지 #1~#3과 텍스트 #1~#3이 같은 번호로 대응해 각 수정의 참조 범위가 고정됐다.", success=True,
    )
    add_pattern_example_slide(
        prs, 5, "이미지가 성공한 방식 3", "삭제 대상과 보존 대상의 경계를 명시",
        [IMAGES_DIR / "064.jpg"],
        "스크린샷처럼 시작일 설정과 날짜별 기록 수정 두 버튼만 남기고, 출석체크 90의 큰 박스와 내용물은 제거해줘.",
        "P01 · 출석체크90 136행", "성공",
        "'두 버튼만 남김'과 '큰 박스·내용물 제거'를 함께 적어 삭제 경계와 보존 경계를 동시에 확정했다.", success=True,
    )
    add_pattern_example_slide(
        prs, 6, "이미지가 실패한 방식 1", "'이렇게'라는 말로 이미지 해석에 전적으로 의존",
        [IMAGES_DIR / "049.jpg"],
        "이렇게 수정. 사진 이해 못했으면 말해.",
        "P01 · 영양분석 108행", "실패",
        "수정할 대상·동작·완성 상태가 없어 AI가 이미지의 모든 차이를 스스로 추측해야 했다.", success=False,
    )
    add_pattern_example_slide(
        prs, 7, "이미지가 실패한 방식 2", "'다시'라는 말만 추가하고 이전 실패의 원인·변경점을 생략",
        [IMAGES_DIR / "050.jpg"],
        "다시. 화나게 하지 말고.",
        "P01 · 영양분석 110행", "실패",
        "이전 요청과 달라진 정보가 없어 실패한 해석을 교정할 단서가 제공되지 않았다.", success=False,
    )
    add_pattern_example_slide(
        prs, 8, "이미지가 실패한 방식 3", "부모 컨테이너와 내부 요소의 삭제 경계가 불분명",
        [IMAGES_DIR / "063.jpg"],
        "출석체크 90 박스 안의 내용을 스크린샷처럼 변경해줘.",
        "P01 · 출석체크90 135행", "일부 실패",
        "박스·내용·버튼 중 무엇을 지우고 남길지 명시되지 않아 요청하지 않은 텍스트가 추가됐다. 136행에서 삭제/보존 경계를 적은 후 성공했다.", success=False,
    )

    add_implication_slide(
        prs, 9, 1, "터치로 수정 대상 선택", "Flutter 컴포넌트 ID·Semantics 기반 오버레이",
        "10/14\n71.4%", "시각 표식+명시적 UI 동작 엄격 성공",
        ["수정할 UI 위에 원·색·블록을 직접 표시", "삭제·이동·배치 텍스트를 함께 제공", "일부 실패 포함 13/14(92.9%)"],
        "빨간 동그라미의 년월을 달력 상단 중앙에… 화살표는 년월 왼쪽·오른쪽에 배치해줘.", "P01 · 대화일정 144행 · 성공",
        "사용자가 터치한 UI를 하이라이트하고 component_id·부모 관계·현재 속성을 Codex에 전달.", image=IMAGES_DIR / "068.jpg",
    )
    add_implication_slide(
        prs, 10, 2, "표준 심볼 스크린샷 편집기", "사각형=대상, 화살표=이동, X=삭제, 번호=요구사항",
        "4/9\n44.4%", "임의 번호 매핑 엄격 성공",
        ["1·2·3, A/B, 색상의 의미가 요청마다 다름", "일부 실패 포함 7/9(77.8%)", "임의 심볼은 하위 요구 누락을 만듦"],
        "2번 블록을 메인화면에서 지워. 지운 2번 블록은 1번 아이콘을 누르면 따로 띄우게…", "P01 · 금융브리핑 42행 · 성공",
        "심볼을 구조화된 대상·동작·목표 위치로 변환하고, 전송 전 사용자에게 확인.", image=IMAGES_DIR / "015.jpg", color=BLUE,
    )
    add_implication_slide(
        prs, 11, 3, "수정 요청 컴파일러", "대상·동작·기대 결과·보존 항목으로 변환",
        "87.5%\nvs 0%", "구체 배치 14/16 vs 모호 이미지 0/6",
        ["완성 상태·순서·보존 대상을 적은 요청이 높은 성공", "'이렇게', '수정해줘'는 엄격 성공 0/6", "빈 정보만 묻는 보조 절차가 필요"],
        "레시피 검색을 맨 위로… 추천·즐겨찾기·내정보는 가로 한 줄로 맨 밑에… 스크롤이 아니라 화면 전환으로…", "P01 · 맞춤레시피 59행 · 성공",
        "AI가 자유 입력에서 네 항목을 추출하고 빈 항목만 물은 후, 확인된 요청을 DB와 Codex에 동일하게 저장.", image=IMAGES_DIR / "026.jpg",
    )
    add_implication_slide(
        prs, 12, 4, "이미지만 전송 시 의도 확인", "추정한 대상·동작을 확인한 후 수정 시작",
        "4/5\n80.0%", "현재 라벨의 '성공'; 주로 질문 생성",
        ["입력 원문은 '(이미지만 첨부)'", "성공 라벨은 앱 수정이 아닌 확인 질문인 경우가 많음", "추측한 의도로 바로 수정하면 오판 가능"],
        "(이미지만 첨부)", "P01 · 원본 95·115·117·119·121행",
        "AI가 추정한 대상·동작·보존 요소를 1~3개 질문으로 제시하고 사용자 확인 후만 Codex에 전달.", image=IMAGES_DIR / "055.jpg", color=BLUE,
    )
    add_implication_slide(
        prs, 13, 5, "UI/시스템 문제 자동 분기", "스크린샷을 보낼 문제와 로그·상태를 보낼 문제 분리",
        "0/5\nvs 1/1", "시스템 이미지 vs 수용 기준 구조화",
        ["권한·위젯 스크린샷 5건 모두 실패", "권한 갱신·서비스 연결·수집 조건을 적은 146행 성공", "화면은 현상이지 원인이 아님"],
        "권한 허용을 하여도 알림 접근 화면에 변화가 일어나지 않아… 체크 표시를 해줬으면 좋겠어.", "P01 · 대화일정 141행 · 실패",
        "UI는 시각 태깅, 권한·서비스·위젯은 매니페스트·서비스 상태·logcat·재현 절차를 자동 수집.", image=IMAGES_DIR / "067.jpg", color=ORANGE,
    )
    add_implication_slide(
        prs, 14, 6, "복수 변경을 항목별로 분할", "원자적 작업·항목별 성공 기준·실패 항목만 재시도",
        "5/18\n27.8%", "번호화된 복수 요청의 미완전 성공",
        ["복수 요청 18건 중 13건은 성공", "5건은 일부 실패 또는 실패", "하나의 라벨로는 어느 항목이 누락됐는지 파악 불가"],
        "1. 1번 블록은 삭제  2. 2번 블록은 간소화  3. 3번 블록에서 빨간색으로 표시한 부분 삭제", "P01 · 금융브리핑 43행 · 일부 실패",
        "번호별로 완료·실패·보류를 저장하고, 실패한 항목만 새 리비전에서 재시도.", image=IMAGES_DIR / "016.jpg", color=ORANGE,
    )
    add_implication_slide(
        prs, 15, 7, "실패 시 범위 축소·폴백", "핵심 기능의 최소 성공 범위를 먼저 만든 후 확장",
        "1/4\n25.0%", "대화일정 위젯 연속 시도 성공",
        ["오늘·내일 일정 위젯 3회 실패", "마지막 요청에서 규격 축소", "불가능하면 '오늘 일정만'이라는 폴백을 허용한 후 성공"],
        "위젯 규격을 줄여서 다시 만들어줘. 정보가 너무 많아 불가능하다면 오늘 일정만 보이게 만들어줘도 좋아.", "P01 · 대화일정 154행 · 성공",
        "실패 시 2~3개의 축소 범위를 제안하고, 사용자가 허용한 폴백을 요청 메타데이터로 저장.", image=IMAGES_DIR / "073.jpg",
    )
    add_implication_slide(
        prs, 16, 8, "런타임 오류 증거 번들", "문제 유형별로 화면·상호작용·로그·권한 상태를 자동 첨부",
        "1/1\nvs 0/4", "UI 오류 증거 vs 시스템 스크린샷",
        ["UI 오류 화면+재현 동작 6행은 성공", "권한·위젯 오류 이미지 4건은 모두 실패", "시스템 문제에는 화면 밖 상태가 필요"],
        "정확한 오류 분석을 위해 오류 화면을 캡처해서 보낼게. 이 오류를 분석하고 다시 오류가 나지 않게 수정해줘.", "P01 · 지도 6행 · 성공",
        "오류 시점의 현재 화면·입력 순서·logcat·권한/서비스 상태·앱 버전을 개인정보 제거 후 첨부.", image=IMAGES_DIR / "001.jpg", color=BLUE,
    )
    add_implication_slide(
        prs, 17, 9, "리비전·보존 대상 자동 첨부", "'아까'·'기존처럼'을 구체적인 버전·컴포넌트로 변환",
        "1/5\n20.0%", "금융브리핑 43~47행 엄격 성공",
        ["'아까', 'v2에서 다시', '메인화면엔 그대로' 사용", "5건 중 4건은 일부 실패", "코드 현재 상태와 참가자가 기억한 상태가 다를 수 있음"],
        "금융브리핑 v2에서 다시 시작할게. 아까 삭제했던 부분에… 메인화면엔 냅두고…", "P01 · 금융브리핑 44행 · 일부 실패",
        "기준 리비전, 선택한 component_id, 보존할 기능을 수정 요청에 자동 첨부.", image=IMAGES_DIR / "017.jpg", color=ORANGE,
    )
    add_implication_slide(
        prs, 18, 10, "성공률을 문제 유형·기술 실패와 분리", "매체 효과를 난이도·빌드·설치·런타임 성공과 함께 층화",
        "65.6%\nvs 74.3%", "이미지 42/64 vs 텍스트 55/74",
        ["이미지는 복잡하거나 이미 실패한 문제에 사용되는 경향", "146행과 150행은 거의 동일한 요청", "150행은 표현 문제가 아니라 Flutter build 실패"],
        "권한 허용 후 공식 Android API로 다시 확인해 1초 이내 체크 표시를 갱신해줘… 서비스 연결 후에만 활성 알림을 조회해줘.", "146행 성공 / 150행 동일 요청·빌드 실패",
        "UI 이동·상태 로직·외부 연동·위젯으로 층화하고, 기술 실패를 '사용자 표현 실패'와 별도 측정.", image=IMAGES_DIR / "070.jpg", color=BLUE,
    )
    add_implication_slide(
        prs, 19, 11, "이미지만 첨부한 성공을 이중 정의", "'적절한 확인 질문'과 '최종 앱 반영'을 별도 라벨로 저장",
        "4/5\n80.0%", "현재 이미지 단독 성공 라벨",
        ["현재 성공은 주로 '질문을 잘 생성함'을 의미", "질문 성공은 앱 수정 성공이 아님", "하나의 라벨로 합치면 이미지 효과가 과대 평가"],
        "(이미지만 첨부) → AI가 화면에서 추정한 수정 대상과 의도를 질문", "P01 · 산부인과 용어집 등 5건",
        "DB에 '확인 질문 적절성'·'사용자 확인 여부'·'앱 수정 반영 여부'를 별도 컬럼으로 저장.", image=IMAGES_DIR / "060.jpg",
    )

    slide = add_new_slide(prs, "결론과 구현 우선순위", 20)
    add_text(slide, 0.6, 1.18, 8.8, 0.55, "이미지 첨부 자체보다 의도를 구조화하는 상호작용이 핵심", size=23, bold=True)
    priorities = [
        ("1", "컴포넌트 선택 오버레이 + 표준 심볼 편집기", "UI의 '어디'를 직접 지정"),
        ("2", "대상·동작·기대 결과·보존 요청 컴파일러", "모호한 입력을 검증 가능한 요청으로 변환"),
        ("3", "UI/시스템 문제 분기 + 런타임 증거", "권한·위젯·외부 연동 진단 근거 확보"),
        ("4", "항목별 성공·기술 실패 분리 측정", "기능의 효과를 왜곡 없이 검증"),
    ]
    for index, (number, heading, body) in enumerate(priorities):
        y = 1.95 + index * 0.76
        color = TEAL if index < 2 else BLUE
        add_rect(slide, 0.7, y, 0.48, 0.48, color, radius=True)
        add_text(slide, 0.71, y + 0.08, 0.46, 0.25, number, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, 1.4, y - 0.02, 5.15, 0.31, heading, size=15.5, bold=True)
        add_text(slide, 6.35, y - 0.01, 2.9, 0.45, body, size=11.3, color=MUTED)

    prs.save(OUTPUT_PPTX)


def validate(rows: dict[int, dict[str, str]]) -> None:
    presentation = Presentation(OUTPUT_PPTX)
    assert len(presentation.slides) == 20
    for source_row in (141, 144, 146, 154):
        assert source_row in rows
    workbook = load_workbook(WORKBOOK_PATH, read_only=True, data_only=False)
    implication_ws = workbook["Implication"]
    values = [clean(implication_ws.cell(row, 4).value) for row in range(2, implication_ws.max_row + 1)]
    assert any("71.4%" in value for value in values)
    assert any("0건(0%)" in value for value in values)


def main() -> None:
    workbook = load_workbook(WORKBOOK_PATH)
    rows = load_labeled_rows(workbook)
    update_workbook_evidence(workbook, rows)
    workbook.save(WORKBOOK_PATH)
    build_deck(rows)
    validate(rows)
    print(f"Updated workbook evidence: {WORKBOOK_PATH}")
    print(f"Updated evidence deck: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
